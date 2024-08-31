from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Inches as DocxInches
from docx.enum.text import WD_ALIGN_PARAGRAPH
import io
import dash
from dash import dcc, html, Input, Output, State
import dash_bootstrap_components as dbc
import pandas as pd
import plotly.express as px
import plotly.graph_objs as go
import numpy as np
from collections import Counter
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize

def generate_ppt_and_doc(movie_lines, movie_conversations):
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    main_color = RGBColor(58, 80, 107)
    accent_color = RGBColor(28, 37, 65)
    background_color = RGBColor(255, 255, 255)

    def add_title_slide(title, subtitle):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = background_color
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        title_shape.text = title
        subtitle_shape.text = subtitle
        title_shape.text_frame.paragraphs[0].font.color.rgb = main_color
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = accent_color

    def add_content_slide(title, content):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = background_color
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        title_shape.text = title
        content_shape.text = content
        title_shape.text_frame.paragraphs[0].font.color.rgb = main_color
        for paragraph in content_shape.text_frame.paragraphs:
            paragraph.font.color.rgb = accent_color

    add_title_slide("Movie Dialogue Analysis Dashboard", "Created by Pallavi Sikha")

    add_content_slide("Dataset Overview",
                      f"Total number of lines: {len(movie_lines)}\n"
                      f"Total number of conversations: {len(movie_conversations)}\n"
                      f"Number of unique characters: {len(set(movie_lines['characterID']))}\n"
                      f"Number of unique movies: {len(set(movie_lines['movieID']))}")

    add_content_slide("Top Characters",
                      "Top 5 characters by number of lines:\n" +
                      "\n".join([f"{i+1}. Character {char}" for i, char in
                                 enumerate(movie_lines['characterID'].value_counts().nlargest(5).index)]))

    add_content_slide("Dialogue Statistics",
                      f"Average words per line: {movie_lines['text'].str.split().str.len().mean():.2f}\n"
                      f"Longest line (words): {movie_lines['text'].str.split().str.len().max()}\n"
                      f"Shortest line (words): {movie_lines['text'].str.split().str.len().min()}")

    add_content_slide("Conversation Analysis",
                      f"Average turns per conversation: {movie_conversations['utteranceIDs'].apply(lambda x: len(eval(x))).mean():.2f}\n"
                      f"Longest conversation (turns): {movie_conversations['utteranceIDs'].apply(lambda x: len(eval(x))).max()}\n"
                      f"Shortest conversation (turns): {movie_conversations['utteranceIDs'].apply(lambda x: len(eval(x))).min()}")

    add_content_slide("Next Steps",
                      "1. Perform sentiment analysis on dialogues\n"
                      "2. Analyze character relationships\n"
                      "3. Study dialogue patterns across different movies\n"
                      "4. Implement advanced NLP techniques for deeper insights\n"
                      "5. Develop a recommendation system based on dialogue preferences")

    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)

    doc = Document()

    def add_title(title):
        doc.add_heading(title, level=1)

    def add_paragraph(text):
        p = doc.add_paragraph(text)
        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

    doc.add_heading('Movie Dialogue Analysis Report', 0)
    doc.add_paragraph('Created by Pallavi Sikha')

    add_title("1. Introduction")
    add_paragraph("This report presents an analysis of movie dialogues based on the Cornell Movie Dialogs Corpus. "
                  "The analysis covers various aspects of the dataset, including dialogue statistics, character interactions, "
                  "and conversation patterns.")

    add_title("2. Dataset Overview")
    add_paragraph(f"The dataset contains {len(movie_lines)} lines of dialogue from {len(set(movie_lines['movieID']))} unique movies. "
                  f"There are {len(set(movie_lines['characterID']))} distinct characters engaged in {len(movie_conversations)} conversations.")

    add_title("3. Dialogue Statistics")
    add_paragraph(f"On average, each line of dialogue contains {movie_lines['text'].str.split().str.len().mean():.2f} words. "
                  f"The longest line has {movie_lines['text'].str.split().str.len().max()} words, while the shortest has "
                  f"{movie_lines['text'].str.split().str.len().min()} words.")

    add_title("4. Character Analysis")
    add_paragraph("The top 5 characters with the most lines of dialogue are:")
    for i, char in enumerate(movie_lines['characterID'].value_counts().nlargest(5).index):
        doc.add_paragraph(f"{i+1}. Character {char}", style='List Bullet')

    add_title("5. Conversation Patterns")
    add_paragraph(f"Conversations in the dataset have an average of {movie_conversations['utteranceIDs'].apply(lambda x: len(eval(x))).mean():.2f} turns. "
                  f"The longest conversation consists of {movie_conversations['utteranceIDs'].apply(lambda x: len(eval(x))).max()} turns, "
                  f"while the shortest has {movie_conversations['utteranceIDs'].apply(lambda x: len(eval(x))).min()} turns.")

    add_title("6. Conclusion and Next Steps")
    add_paragraph("This analysis provides initial insights into the movie dialogue dataset. To further enhance our understanding, "
                  "we propose the following next steps:")
    doc.add_paragraph("1. Conduct sentiment analysis on the dialogues", style='List Number')
    doc.add_paragraph("2. Analyze character relationships and interaction patterns", style='List Number')
    doc.add_paragraph("3. Study dialogue variations across different movie genres", style='List Number')
    doc.add_paragraph("4. Implement advanced NLP techniques for deeper content analysis", style='List Number')
    doc.add_paragraph("5. Develop a recommendation system based on dialogue preferences", style='List Number')

    # Save Word document
    doc_buffer = io.BytesIO()
    doc.save(doc_buffer)
    doc_buffer.seek(0)

    return ppt_buffer, doc_buffer

@app.callback(
    Output('download-ppt', 'data'),
    Input('generate-ppt-button', 'n_clicks'),
    prevent_initial_call=True
)
def download_ppt(n_clicks):
    ppt_buffer, _ = generate_ppt_and_doc(movie_lines, movie_conversations)
    return dcc.send_bytes(ppt_buffer.getvalue(), "movie_dialogue_analysis.pptx")

@app.callback(
    Output('download-doc', 'data'),
    Input('generate-doc-button', 'n_clicks'),
    prevent_initial_call=True
)
def download_doc(n_clicks):
    _, doc_buffer = generate_ppt_and_doc(movie_lines, movie_conversations)
    return dcc.send_bytes(doc_buffer.getvalue(), "movie_dialogue_analysis.docx")